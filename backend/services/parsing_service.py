from typing import Dict, List, Optional, Union, Generator, Set, Tuple
from enum import Enum
import logging
from pathlib import Path
import PyPDF2
from docx import Document
import pandas as pd
from pydantic import BaseModel
import markdown
from bs4 import BeautifulSoup
import openpyxl
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from langdetect import detect, LangDetectException
# import pytesseract
# from PIL import Image
import io
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
import asyncio
from datetime import datetime
import traceback
from typing import Type, Dict, Any

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ParsingError(Exception):
    """解析错误基类"""
    def __init__(self, message: str, file_path: str, error_code: str = "UNKNOWN_ERROR"):
        self.message = message
        self.file_path = file_path
        self.error_code = error_code
        super().__init__(self.message)

class FileNotFoundError(ParsingError):
    """文件不存在错误"""
    def __init__(self, file_path: str):
        super().__init__(f"文件不存在: {file_path}", file_path, "FILE_NOT_FOUND")

class UnsupportedFileTypeError(ParsingError):
    """不支持的文件类型错误"""
    def __init__(self, file_path: str, file_type: str):
        super().__init__(
            f"不支持的文件类型: {file_type}",
            file_path,
            "UNSUPPORTED_FILE_TYPE"
        )

class FileReadError(ParsingError):
    """文件读取错误"""
    def __init__(self, file_path: str, error: Exception):
        super().__init__(
            f"文件读取失败: {str(error)}",
            file_path,
            "FILE_READ_ERROR"
        )

class FileParseError(ParsingError):
    """文件解析错误"""
    def __init__(self, file_path: str, error: Exception):
        super().__init__(
            f"文件解析失败: {str(error)}",
            file_path,
            "FILE_PARSE_ERROR"
        )

class MetadataExtractionError(ParsingError):
    """元数据提取错误"""
    def __init__(self, file_path: str, error: Exception):
        super().__init__(
            f"元数据提取失败: {str(error)}",
            file_path,
            "METADATA_EXTRACTION_ERROR"
        )

class ChunkProcessingError(ParsingError):
    """分片处理错误"""
    def __init__(self, file_path: str, chunk_index: int, error: Exception):
        super().__init__(
            f"分片处理失败 [chunk {chunk_index}]: {str(error)}",
            file_path,
            "CHUNK_PROCESSING_ERROR"
        )

class BatchProcessingError(ParsingError):
    """批量处理错误"""
    def __init__(self, file_path: str, error: Exception):
        super().__init__(
            f"批量处理失败: {str(error)}",
            file_path,
            "BATCH_PROCESSING_ERROR"
        )

class ProcessingTimeoutError(ParsingError):
    """处理超时错误"""
    def __init__(self, file_path: str, timeout: float):
        super().__init__(
            f"处理超时 (超过 {timeout} 秒)",
            file_path,
            "PROCESSING_TIMEOUT"
        )

class ProcessingCancelledError(ParsingError):
    """处理取消错误"""
    def __init__(self, file_path: str):
        super().__init__(
            "处理被取消",
            file_path,
            "PROCESSING_CANCELLED"
        )

class ErrorHandler:
    """错误处理器"""
    
    def __init__(self):
        self.error_handlers: Dict[str, callable] = {
            "FILE_NOT_FOUND": self._handle_file_not_found,
            "UNSUPPORTED_FILE_TYPE": self._handle_unsupported_file_type,
            "FILE_READ_ERROR": self._handle_file_read_error,
            "FILE_PARSE_ERROR": self._handle_file_parse_error,
            "METADATA_EXTRACTION_ERROR": self._handle_metadata_extraction_error,
            "CHUNK_PROCESSING_ERROR": self._handle_chunk_processing_error,
            "BATCH_PROCESSING_ERROR": self._handle_batch_processing_error,
            "PROCESSING_TIMEOUT": self._handle_processing_timeout,
            "PROCESSING_CANCELLED": self._handle_processing_cancelled,
            "UNKNOWN_ERROR": self._handle_unknown_error
        }
    
    def handle_error(self, error: ParsingError) -> Dict[str, Any]:
        """
        处理错误
        
        Args:
            error: 解析错误
            
        Returns:
            Dict[str, Any]: 错误处理结果
        """
        handler = self.error_handlers.get(error.error_code, self._handle_unknown_error)
        return handler(error)
    
    def _handle_file_not_found(self, error: ParsingError) -> Dict[str, Any]:
        """处理文件不存在错误"""
        logger.error(f"文件不存在: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请检查文件路径是否正确"
        }
    
    def _handle_unsupported_file_type(self, error: ParsingError) -> Dict[str, Any]:
        """处理不支持的文件类型错误"""
        logger.error(f"不支持的文件类型: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请使用支持的文件类型"
        }
    
    def _handle_file_read_error(self, error: ParsingError) -> Dict[str, Any]:
        """处理文件读取错误"""
        logger.error(f"文件读取失败: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请检查文件权限和完整性"
        }
    
    def _handle_file_parse_error(self, error: ParsingError) -> Dict[str, Any]:
        """处理文件解析错误"""
        logger.error(f"文件解析失败: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请检查文件格式是否正确"
        }
    
    def _handle_metadata_extraction_error(self, error: ParsingError) -> Dict[str, Any]:
        """处理元数据提取错误"""
        logger.error(f"元数据提取失败: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请检查文件是否包含有效的元数据"
        }
    
    def _handle_chunk_processing_error(self, error: ParsingError) -> Dict[str, Any]:
        """处理分片处理错误"""
        logger.error(f"分片处理失败: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请尝试减小分片大小或检查文件完整性"
        }
    
    def _handle_batch_processing_error(self, error: ParsingError) -> Dict[str, Any]:
        """处理批量处理错误"""
        logger.error(f"批量处理失败: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请检查系统资源或减小并发数"
        }
    
    def _handle_processing_timeout(self, error: ParsingError) -> Dict[str, Any]:
        """处理处理超时错误"""
        logger.error(f"处理超时: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请增加超时时间或优化处理逻辑"
        }
    
    def _handle_processing_cancelled(self, error: ParsingError) -> Dict[str, Any]:
        """处理处理取消错误"""
        logger.warning(f"处理被取消: {error.file_path}")
        return {
            "status": "cancelled",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "处理已被用户取消"
        }
    
    def _handle_unknown_error(self, error: ParsingError) -> Dict[str, Any]:
        """处理未知错误"""
        logger.error(f"未知错误: {error.file_path}")
        return {
            "status": "error",
            "error_code": error.error_code,
            "message": error.message,
            "file_path": error.file_path,
            "suggestion": "请检查日志获取详细信息"
        }

class ProcessingStatus(str, Enum):
    """处理状态枚举"""
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"

class DocumentMetadata(BaseModel):
    """文档元数据模型"""
    filename: str
    file_type: str
    page_count: Optional[int] = None
    title: Optional[str] = None
    author: Optional[str] = None
    created_date: Optional[str] = None
    modified_date: Optional[str] = None
    language: Optional[str] = None
    word_count: Optional[int] = None
    image_count: Optional[int] = None
    table_count: Optional[int] = None
    file_size: Optional[int] = None
    chunk_count: Optional[int] = None

class ParsedContent(BaseModel):
    """解析后的内容模型"""
    content: str
    metadata: Dict
    page_number: Optional[int] = None
    section_title: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    images: Optional[List[Dict]] = None
    chunk_index: Optional[int] = None
    total_chunks: Optional[int] = None

class ProcessingResult(BaseModel):
    """处理结果模型"""
    file_path: str
    status: ProcessingStatus
    start_time: datetime
    end_time: Optional[datetime] = None
    error: Optional[str] = None
    contents: Optional[List[ParsedContent]] = None
    metadata: Optional[DocumentMetadata] = None

class BatchProcessingConfig(BaseModel):
    """批量处理配置"""
    max_concurrent_tasks: int = 5  # 最大并发任务数
    chunk_size: int = 1024 * 1024  # 分片大小（默认1MB）
    retry_count: int = 3  # 重试次数
    retry_delay: float = 1.0  # 重试延迟（秒）
    timeout: Optional[float] = None  # 超时时间（秒）

class ParsingStrategy(str, Enum):
    """解析策略枚举"""
    ALL_TEXT = "all_text"
    BY_PAGES = "by_pages"
    BY_TITLES = "by_titles"
    TEXT_AND_TABLES = "text_and_tables"

class ParsingService:
    """文档解析服务"""
    
    def __init__(self, config: Optional[BatchProcessingConfig] = None):
        self.supported_extensions = {
            '.pdf', '.docx', '.txt', '.md', '.html', 
            '.xlsx', '.xls', '.pptx', '.rtf'
            # '.png', '.jpg', '.jpeg', '.bmp', '.tiff'  # 暂时不支持图片格式
        }
        self.config = config or BatchProcessingConfig()
        self.max_workers = min(32, (os.cpu_count() or 1) + 4)  # 默认最大线程数
        self._processing_tasks: Dict[str, asyncio.Task] = {}
        self._processing_results: Dict[str, ProcessingResult] = {}
        self.error_handler = ErrorHandler()
    
    def _get_file_chunks(self, file_path: Path) -> Generator[bytes, None, None]:
        """
        将文件分片读取
        
        Args:
            file_path: 文件路径
            
        Yields:
            bytes: 文件块
        """
        with open(file_path, 'rb') as f:
            while True:
                chunk = f.read(self.config.chunk_size)
                if not chunk:
                    break
                yield chunk
    
    async def _process_chunk(self, chunk: bytes, chunk_index: int, total_chunks: int, 
                           file_path: Path, metadata: DocumentMetadata) -> ParsedContent:
        """
        处理单个文件块
        
        Args:
            chunk: 文件块
            chunk_index: 块索引
            total_chunks: 总块数
            file_path: 文件路径
            metadata: 文档元数据
            
        Returns:
            ParsedContent: 解析后的内容
        """
        try:
            # 根据文件类型选择解析方法
            if file_path.suffix.lower() == '.pdf':
                content = self._process_pdf_chunk(chunk)
            elif file_path.suffix.lower() in {'.docx', '.doc'}:
                content = self._process_docx_chunk(chunk)
            elif file_path.suffix.lower() in {'.txt', '.md', '.html', '.rtf'}:
                content = self._process_text_chunk(chunk)
            else:
                content = chunk.decode('utf-8', errors='ignore')
            
            return ParsedContent(
                content=content,
                metadata=metadata.dict(),
                chunk_index=chunk_index,
                total_chunks=total_chunks
            )
        except Exception as e:
            logger.error(f"处理文件块失败: {str(e)}")
            raise
    
    def _process_pdf_chunk(self, chunk: bytes) -> str:
        """处理PDF文件块"""
        try:
            pdf_file = io.BytesIO(chunk)
            pdf = PyPDF2.PdfReader(pdf_file)
            return " ".join([page.extract_text() for page in pdf.pages])
        except Exception as e:
            logger.error(f"PDF块处理失败: {str(e)}")
            return ""
    
    def _process_docx_chunk(self, chunk: bytes) -> str:
        """处理Word文件块"""
        try:
            doc_file = io.BytesIO(chunk)
            doc = Document(doc_file)
            return " ".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"Word块处理失败: {str(e)}")
            return ""
    
    def _process_text_chunk(self, chunk: bytes) -> str:
        """处理文本文件块"""
        try:
            return chunk.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"文本块处理失败: {str(e)}")
            return ""
    
    async def parse_document(
        self,
        file_path: Union[str, Path],
        strategy: ParsingStrategy = ParsingStrategy.ALL_TEXT,
        **kwargs
    ) -> List[ParsedContent]:
        """
        解析文档
        
        Args:
            file_path: 文档路径
            strategy: 解析策略
            **kwargs: 其他参数
            
        Returns:
            List[ParsedContent]: 解析后的内容列表
            
        Raises:
            ParsingError: 解析过程中的各种错误
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(str(file_path))
            
            if file_path.suffix.lower() not in self.supported_extensions:
                raise UnsupportedFileTypeError(str(file_path), file_path.suffix)
            
            # 获取文件大小
            try:
                file_size = file_path.stat().st_size
            except Exception as e:
                raise FileReadError(str(file_path), e)
            
            # 如果文件小于分片大小，直接处理
            if file_size <= self.config.chunk_size:
                return await self._parse_small_file(file_path, strategy, **kwargs)
            
            # 提取文档元数据
            try:
                metadata = await self._extract_metadata(file_path)
            except Exception as e:
                raise MetadataExtractionError(str(file_path), e)
            
            metadata.file_size = file_size
            
            # 计算分片数量
            total_chunks = (file_size + self.config.chunk_size - 1) // self.config.chunk_size
            metadata.chunk_count = total_chunks
            
            # 使用线程池并行处理文件块
            contents = []
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                futures = []
                for i, chunk in enumerate(self._get_file_chunks(file_path)):
                    future = executor.submit(
                        self._process_chunk,
                        chunk,
                        i,
                        total_chunks,
                        file_path,
                        metadata
                    )
                    futures.append(future)
                
                # 收集处理结果
                for future in as_completed(futures):
                    try:
                        content = future.result()
                        contents.append(content)
                    except Exception as e:
                        raise ChunkProcessingError(str(file_path), i, e)
            
            # 按块索引排序
            contents.sort(key=lambda x: x.chunk_index)
            
            return contents
                
        except ParsingError:
            raise
        except Exception as e:
            logger.error(f"文档解析失败: {str(e)}\n{traceback.format_exc()}")
            raise FileParseError(str(file_path), e)
    
    async def _parse_small_file(
        self,
        file_path: Path,
        strategy: ParsingStrategy,
        **kwargs
    ) -> List[ParsedContent]:
        """处理小文件"""
        try:
            # 提取文档元数据
            metadata = await self._extract_metadata(file_path)
            
            # 根据文件类型选择解析方法
            if file_path.suffix.lower() in {'.xlsx', '.xls'}:
                return await self._parse_excel(file_path, metadata)
            elif file_path.suffix.lower() == '.pptx':
                return await self._parse_powerpoint(file_path, metadata)
            elif file_path.suffix.lower() == '.rtf':
                return await self._parse_rtf(file_path, metadata)
            elif file_path.suffix.lower() == '.md':
                return await self._parse_markdown(file_path, metadata)
            elif file_path.suffix.lower() == '.html':
                return await self._parse_html(file_path, metadata)
            
            # 根据策略解析文档
            if strategy == ParsingStrategy.ALL_TEXT:
                return await self._parse_all_text(file_path, metadata)
            elif strategy == ParsingStrategy.BY_PAGES:
                return await self._parse_by_pages(file_path, metadata)
            elif strategy == ParsingStrategy.BY_TITLES:
                return await self._parse_by_titles(file_path, metadata)
            elif strategy == ParsingStrategy.TEXT_AND_TABLES:
                return await self._parse_text_and_tables(file_path, metadata)
            else:
                raise ValueError(f"不支持的解析策略: {strategy}")
        except Exception as e:
            logger.error(f"小文件解析失败: {str(e)}")
            raise

    # async def _parse_image(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
    #     """解析图片文件"""
    #     try:
    #         with Image.open(file_path) as img:
    #             # 使用 OCR 提取文本
    #             text = pytesseract.image_to_string(img)
    #             
    #             # 更新元数据
    #             metadata.image_count = 1
    #             metadata.word_count = len(text.split())
    #             
    #             return [ParsedContent(
    #                 content=text,
    #                 metadata=metadata.dict(),
    #                 images=[{
    #                     'format': img.format,
    #                     'size': img.size,
    #                     'mode': img.mode
    #                 }]
    #             )]
    #     except Exception as e:
    #         logger.error(f"图片解析失败: {str(e)}")
    #         raise

    async def _parse_excel(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """解析 Excel 文件"""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            contents = []
            
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                # 将工作表转换为 DataFrame
                data = ws.values
                cols = next(data)
                df = pd.DataFrame(data, columns=cols)
                
                # 提取文本内容
                text = "\n".join([str(cell) for row in ws.rows for cell in row if cell.value])
                
                contents.append(ParsedContent(
                    content=text,
                    metadata=metadata.dict(),
                    tables=[df],
                    section_title=sheet
                ))
            
            # 更新元数据
            metadata.table_count = len(contents)
            metadata.word_count = sum(len(c.content.split()) for c in contents)
            
            return contents
        except Exception as e:
            logger.error(f"Excel 解析失败: {str(e)}")
            raise

    async def _parse_powerpoint(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """解析 PowerPoint 文件"""
        try:
            prs = Presentation(file_path)
            contents = []
            
            for i, slide in enumerate(prs.slides, 1):
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                
                contents.append(ParsedContent(
                    content="\n".join(text),
                    metadata=metadata.dict(),
                    page_number=i
                ))
            
            # 更新元数据
            metadata.page_count = len(contents)
            metadata.word_count = sum(len(c.content.split()) for c in contents)
            
            return contents
        except Exception as e:
            logger.error(f"PowerPoint 解析失败: {str(e)}")
            raise

    async def _parse_rtf(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """解析 RTF 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_text = f.read()
            
            # 转换 RTF 为纯文本
            text = rtf_to_text(rtf_text)
            
            # 更新元数据
            metadata.word_count = len(text.split())
            
            return [ParsedContent(
                content=text,
                metadata=metadata.dict()
            )]
        except Exception as e:
            logger.error(f"RTF 解析失败: {str(e)}")
            raise

    async def _parse_markdown(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """解析 Markdown 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_text = f.read()
            
            # 转换为 HTML
            html = markdown.markdown(md_text)
            
            # 使用 BeautifulSoup 提取文本
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
            # 更新元数据
            metadata.word_count = len(text.split())
            
            return [ParsedContent(
                content=text,
                metadata=metadata.dict()
            )]
        except Exception as e:
            logger.error(f"Markdown 解析失败: {str(e)}")
            raise

    async def _parse_html(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """解析 HTML 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_text = f.read()
            
            # 使用 BeautifulSoup 解析 HTML
            soup = BeautifulSoup(html_text, 'html.parser')
            
            # 提取文本
            text = soup.get_text()
            
            # 提取表格
            tables = []
            for table in soup.find_all('table'):
                df = pd.read_html(str(table))[0]
                tables.append(df)
            
            # 更新元数据
            metadata.word_count = len(text.split())
            metadata.table_count = len(tables)
            
            return [ParsedContent(
                content=text,
                metadata=metadata.dict(),
                tables=tables
            )]
        except Exception as e:
            logger.error(f"HTML 解析失败: {str(e)}")
            raise

    async def _detect_language(self, text: str) -> str:
        """
        检测文本语言
        
        Args:
            text: 要检测的文本
            
        Returns:
            str: 语言代码（如 'en', 'zh-cn', 'ja' 等）
        """
        try:
            # 如果文本太短，可能无法准确检测
            if len(text.strip()) < 10:
                return "unknown"
            
            # 检测语言
            lang = detect(text)
            return lang
        except LangDetectException as e:
            logger.warning(f"语言检测失败: {str(e)}")
            return "unknown"

    async def _extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """提取文档元数据"""
        try:
            # 读取文件内容用于语言检测
            content = ""
            if file_path.suffix == '.pdf':
                with open(file_path, 'rb') as f:
                    pdf = PyPDF2.PdfReader(f)
                    content = " ".join([page.extract_text() for page in pdf.pages])
                    metadata = DocumentMetadata(
                        filename=file_path.name,
                        file_type='pdf',
                        page_count=len(pdf.pages),
                        title=pdf.metadata.get('/Title', ''),
                        author=pdf.metadata.get('/Author', '')
                    )
            elif file_path.suffix == '.docx':
                doc = Document(file_path)
                content = " ".join([p.text for p in doc.paragraphs])
                core_props = doc.core_properties
                metadata = DocumentMetadata(
                    filename=file_path.name,
                    file_type='docx',
                    title=core_props.title,
                    author=core_props.author,
                    created_date=str(core_props.created),
                    modified_date=str(core_props.modified)
                )
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                metadata = DocumentMetadata(
                    filename=file_path.name,
                    file_type=file_path.suffix[1:]
                )
            
            # 检测语言
            metadata.language = await self._detect_language(content)
            
            return metadata
        except Exception as e:
            logger.error(f"元数据提取失败: {str(e)}")
            return DocumentMetadata(filename=file_path.name, file_type=file_path.suffix[1:])
    
    async def _parse_all_text(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """全文提取策略"""
        try:
            if file_path.suffix == '.pdf':
                with open(file_path, 'rb') as f:
                    pdf = PyPDF2.PdfReader(f)
                    text = ""
                    for page in pdf.pages:
                        text += page.extract_text() + "\n"
                    return [ParsedContent(content=text, metadata=metadata.dict())]
            elif file_path.suffix == '.docx':
                doc = Document(file_path)
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                return [ParsedContent(content=text, metadata=metadata.dict())]
            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                return [ParsedContent(content=text, metadata=metadata.dict())]
        except Exception as e:
            logger.error(f"全文提取失败: {str(e)}")
            raise
    
    async def _parse_by_pages(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """逐页解析策略"""
        try:
            if file_path.suffix == '.pdf':
                with open(file_path, 'rb') as f:
                    pdf = PyPDF2.PdfReader(f)
                    contents = []
                    for i, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        contents.append(ParsedContent(
                            content=text,
                            metadata=metadata.dict(),
                            page_number=i
                        ))
                    return contents
            else:
                # 对于非PDF文件，返回单个内容
                return await self._parse_all_text(file_path, metadata)
        except Exception as e:
            logger.error(f"逐页解析失败: {str(e)}")
            raise
    
    async def _parse_by_titles(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """基于标题的分段策略"""
        try:
            if file_path.suffix == '.docx':
                doc = Document(file_path)
                contents = []
                current_section = []
                current_title = None
                
                for paragraph in doc.paragraphs:
                    if paragraph.style.name.startswith('Heading'):
                        if current_section:
                            contents.append(ParsedContent(
                                content="\n".join(current_section),
                                metadata=metadata.dict(),
                                section_title=current_title
                            ))
                        current_title = paragraph.text
                        current_section = []
                    else:
                        current_section.append(paragraph.text)
                
                if current_section:
                    contents.append(ParsedContent(
                        content="\n".join(current_section),
                        metadata=metadata.dict(),
                        section_title=current_title
                    ))
                return contents
            else:
                # 对于非Word文档，返回单个内容
                return await self._parse_all_text(file_path, metadata)
        except Exception as e:
            logger.error(f"标题分段解析失败: {str(e)}")
            raise
    
    async def _parse_text_and_tables(self, file_path: Path, metadata: DocumentMetadata) -> List[ParsedContent]:
        """文本和表格混合解析策略"""
        try:
            if file_path.suffix == '.docx':
                doc = Document(file_path)
                contents = []
                current_text = []
                
                for element in doc.element.body:
                    if element.tag.endswith('tbl'):
                        # 处理表格
                        table = pd.read_html(element.xml)[0]
                        if current_text:
                            contents.append(ParsedContent(
                                content="\n".join(current_text),
                                metadata=metadata.dict()
                            ))
                            current_text = []
                        contents.append(ParsedContent(
                            content="",
                            metadata=metadata.dict(),
                            tables=[table]
                        ))
                    else:
                        # 处理文本
                        text = element.text
                        if text:
                            current_text.append(text)
                
                if current_text:
                    contents.append(ParsedContent(
                        content="\n".join(current_text),
                        metadata=metadata.dict()
                    ))
                return contents
            else:
                # 对于非Word文档，返回单个内容
                return await self._parse_all_text(file_path, metadata)
        except Exception as e:
            logger.error(f"文本和表格混合解析失败: {str(e)}")
            raise

    async def process_batch(
        self,
        file_paths: List[Union[str, Path]],
        strategy: ParsingStrategy = ParsingStrategy.ALL_TEXT,
        **kwargs
    ) -> Dict[str, ProcessingResult]:
        """
        批量处理文档
        
        Args:
            file_paths: 文件路径列表
            strategy: 解析策略
            **kwargs: 其他参数
            
        Returns:
            Dict[str, ProcessingResult]: 处理结果字典
        """
        try:
            # 初始化处理结果
            for file_path in file_paths:
                file_path = str(file_path)
                if file_path not in self._processing_results:
                    self._processing_results[file_path] = ProcessingResult(
                        file_path=file_path,
                        status=ProcessingStatus.PENDING,
                        start_time=datetime.now()
                    )
            
            # 创建处理任务
            tasks = []
            for file_path in file_paths:
                file_path = str(file_path)
                if self._processing_results[file_path].status == ProcessingStatus.PENDING:
                    task = asyncio.create_task(
                        self._process_single_file(file_path, strategy, **kwargs)
                    )
                    self._processing_tasks[file_path] = task
                    tasks.append(task)
            
            # 使用信号量控制并发
            semaphore = asyncio.Semaphore(self.config.max_concurrent_tasks)
            
            async def process_with_semaphore(task):
                async with semaphore:
                    return await task
            
            # 等待所有任务完成
            results = await asyncio.gather(
                *[process_with_semaphore(task) for task in tasks],
                return_exceptions=True
            )
            
            # 更新处理结果
            for file_path, result in zip(file_paths, results):
                file_path = str(file_path)
                if isinstance(result, Exception):
                    error_result = self.error_handler.handle_error(
                        BatchProcessingError(file_path, result)
                    )
                    self._processing_results[file_path].status = ProcessingStatus.FAILED
                    self._processing_results[file_path].error = error_result["message"]
                else:
                    self._processing_results[file_path].status = ProcessingStatus.COMPLETED
                    self._processing_results[file_path].contents = result
                self._processing_results[file_path].end_time = datetime.now()
            
            return self._processing_results
            
        except Exception as e:
            logger.error(f"批量处理失败: {str(e)}\n{traceback.format_exc()}")
            raise BatchProcessingError(str(file_paths), e)
    
    async def _process_single_file(
        self,
        file_path: str,
        strategy: ParsingStrategy,
        **kwargs
    ) -> List[ParsedContent]:
        """
        处理单个文件（带重试机制）
        
        Args:
            file_path: 文件路径
            strategy: 解析策略
            **kwargs: 其他参数
            
        Returns:
            List[ParsedContent]: 解析后的内容列表
        """
        retry_count = 0
        last_error = None
        
        while retry_count < self.config.retry_count:
            try:
                self._processing_results[file_path].status = ProcessingStatus.PROCESSING
                contents = await self.parse_document(file_path, strategy, **kwargs)
                return contents
            except Exception as e:
                last_error = e
                retry_count += 1
                if retry_count < self.config.retry_count:
                    await asyncio.sleep(self.config.retry_delay)
                    logger.warning(f"重试处理文件 {file_path}，第 {retry_count} 次重试")
                else:
                    logger.error(f"处理文件 {file_path} 失败，已达到最大重试次数")
                    raise last_error
    
    async def get_processing_status(self, file_path: str) -> ProcessingResult:
        """
        获取处理状态
        
        Args:
            file_path: 文件路径
            
        Returns:
            ProcessingResult: 处理结果
        """
        return self._processing_results.get(file_path)
    
    async def cancel_processing(self, file_path: str) -> bool:
        """
        取消处理
        
        Args:
            file_path: 文件路径
            
        Returns:
            bool: 是否成功取消
        """
        if file_path in self._processing_tasks:
            task = self._processing_tasks[file_path]
            if not task.done():
                task.cancel()
                try:
                    await task
                except asyncio.CancelledError:
                    self._processing_results[file_path].status = ProcessingStatus.FAILED
                    self._processing_results[file_path].error = "处理被取消"
                    self._processing_results[file_path].end_time = datetime.now()
                    return True
        return False
    
    async def clear_processing_results(self, file_paths: Optional[List[str]] = None):
        """
        清除处理结果
        
        Args:
            file_paths: 文件路径列表，如果为None则清除所有结果
        """
        if file_paths is None:
            self._processing_results.clear()
            self._processing_tasks.clear()
        else:
            for file_path in file_paths:
                self._processing_results.pop(file_path, None)
                self._processing_tasks.pop(file_path, None) 